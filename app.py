from flask import Flask, request, jsonify, render_template
import sqlite3
import boto3
import chromadb
import requests
import uuid
import os
import json
import pdfplumber
from docx import Document
from pptx import Presentation
import pandas as pd
from urllib.parse import urlparse
from dotenv import load_dotenv
load_dotenv()

app = Flask(__name__)

# Initialize Chroma DB
chroma_client = chromadb.PersistentClient(path='chroma_db')
collection = chroma_client.get_or_create_collection(name="documents")

# SQLite for metadata (Thread-safe connection)
def get_db_connection():
    conn = sqlite3.connect('file_index.db', check_same_thread=False, isolation_level=None)
    conn.row_factory = sqlite3.Row
    return conn

conn = get_db_connection()
c = conn.cursor()
c.execute('''CREATE TABLE IF NOT EXISTS file_indexes
             (clientId INTEGER, fileName TEXT, fileUrl TEXT, propertyId INTEGER, indexedId TEXT)''')
conn.commit()

# Load Ollama API URL from environment variables
OLLAMA_URL = os.getenv("OLLAMA_URL", "http://127.0.0.1:11434/api")


# Load AWS credentials from environment variables (Recommended)
AWS_ACCESS_KEY = os.getenv("AWS_ACCESS_KEY_ID")
AWS_SECRET_KEY = os.getenv("AWS_SECRET_ACCESS_KEY")
S3_BUCKET_NAME = "trojan-horse-ai-qa"
AWS_REGION = "us-west-2"

# Initialize S3 client
s3_client = boto3.client(
    "s3",
    region_name=AWS_REGION,
    aws_access_key_id=AWS_ACCESS_KEY,
    aws_secret_access_key=AWS_SECRET_KEY,
)

def get_file_content(file_path):
    try:
        # Handle S3 files
        if file_path.startswith("s3://"):
            bucket_name, key = file_path.replace("s3://", "").split("/", 1)
            obj = s3_client.get_object(Bucket=bucket_name, Key=key)
            file_content = obj["Body"].read()

            # Auto-detect file type
            if key.endswith(".pdf"):
                return extract_text_from_pdf(file_content)
            elif key.endswith(".docx"):
                return extract_text_from_docx(file_content)
            elif key.endswith(".pptx"):
                return extract_text_from_pptx(file_content)
            elif key.endswith(".xlsx"):
                return extract_text_from_xlsx(file_content)
            else:
                return file_content.decode("utf-8")

        # Handle local file system
        elif os.path.exists(file_path):
            if file_path.endswith(".pdf"):
                return extract_text_from_pdf(file_path)
            elif file_path.endswith(".docx"):
                return extract_text_from_docx(file_path)
            elif file_path.endswith(".pptx"):
                return extract_text_from_pptx(file_path)
            elif file_path.endswith(".xlsx"):
                return extract_text_from_xlsx(file_path)
            else:
                with open(file_path, "r", encoding="utf-8") as f:
                    return f.read()

        return f"Invalid file path: {file_path}"

    except Exception as e:
        return f"Failed to fetch content: {str(e)}"

# Helper functions to extract text
def extract_text_from_pdf(pdf_file):
    text = ""
    with pdfplumber.open(pdf_file) as pdf:
        for page in pdf.pages:
            text += page.extract_text() + "\n"
    return text.strip()

def extract_text_from_docx(docx_file):
    doc = Document(docx_file)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_xlsx(xlsx_file):
    df = pd.read_excel(xlsx_file, engine="openpyxl")
    return df.to_string(index=False)


def get_embedding(text):
    try:
        response = requests.post(
            f"{OLLAMA_URL}/embeddings",
            json={"model": "nomic-embed-text", "prompt": text},
            timeout=30
        )
        response.raise_for_status()
        return response.json().get('embedding', [])
    except Exception as e:
        return f"Embedding failed: {str(e)}"

def generate_response(prompt, model="deepseek-r1:7b"):
    try:
        response = requests.post(
            f"{OLLAMA_URL}/generate",
            json={
                "model": model,
                "prompt": prompt,
                "stream": False,  # Explicitly disable streaming
                "options": {"temperature": 0.3}
            },
            timeout=600
        )
        response.raise_for_status()
        return response.json().get('response', "No response generated")
    except Exception as e:
        return f"LLM failed: {str(e)}"

# Error handlers
@app.errorhandler(400)
def bad_request(e):
    return jsonify(error=str(e)), 400

@app.errorhandler(500)
def server_error(e):
    return jsonify(error=str(e)), 500

# API Endpoints
@app.route('/create-file-indexing', methods=['POST'])
def create_index():
    try:
        data = request.json
        required = ['clientId', 'fileName', 'fileUrl', 'propertyId']
        if not all(k in data for k in required):
            return jsonify(error="Missing required fields"), 400
            
        content = get_file_content(data['fileUrl'])
        embedding = get_embedding(content)
        
        # Generate UUID FIRST before any database operations
        indexed_id = str(uuid.uuid4())  # <--- THIS WAS MISSING/MISPLACED
        
        # Store in Chroma
        collection.add(
            ids=[indexed_id],
            embeddings=[embedding],
            metadatas=[{
                "clientId": int(data['clientId']),
                "propertyId": int(data['propertyId']),
                "fileName": data['fileName'],
                "content": content[:500]
            }],
            documents=[content]
        )
        
        # Store metadata in SQLite
        conn = get_db_connection()
        c = conn.cursor()
        c.execute("INSERT INTO file_indexes VALUES (?, ?, ?, ?, ?)",
                  (data['clientId'], data['fileName'], data['fileUrl'], data['propertyId'], indexed_id))
        conn.commit()
        
        return jsonify({"indexed": True, "indexedid": indexed_id})

    except Exception as e:
        return jsonify(error=str(e)), 500


@app.route('/get-indexed-files', methods=['GET'])
def get_indexed_files():
    try:
        client_id = request.args.get('clientId')
        property_id = request.args.get('propertyId')

        if not client_id or not property_id:
            return jsonify(error="Missing clientId or propertyId"), 400

        conn = get_db_connection()
        c = conn.cursor()
        c.execute("SELECT fileName, fileUrl, indexedId FROM file_indexes WHERE clientId=? AND propertyId=?",
                  (client_id, property_id))
        results = c.fetchall()

        if not results:
            return jsonify(error="No files found for given clientId and propertyId"), 404

        files = [{"fileName": row[0], "fileUrl": row[1], "indexedId": row[2]} for row in results]
        return jsonify(files)

    except Exception as e:
        return jsonify(error=str(e)), 500

@app.route('/get-embedded-docs', methods=['GET'])
def get_embedded_docs():
    try:
        client_id = request.args.get('clientId')
        property_id = request.args.get('propertyId')

        if not client_id or not property_id:
            return jsonify(error="Missing clientId or propertyId"), 400

        # If propertyId is empty or None, return sample-global-faq.txt content
        if not property_id:
            faq_file = "sample-global-faq.txt"
            if not os.path.exists(faq_file):
                return jsonify(error="Global FAQ file not found"), 404

            with open(faq_file, "r", encoding="utf-8") as f:
                faq_content = f.read()

            return jsonify([{"fileName": "sample-global-faq.txt", "contentSnippet": faq_content[:500]}])

        # Convert property_id to int
        property_id = int(property_id)

        # Generate an empty embedding vector to allow querying
        empty_embedding = get_embedding(" ")  # Use space instead of empty string

        # Query ChromaDB using metadata filter and empty embedding
        results = collection.query(
            query_embeddings=[empty_embedding],
            where={
                "$and": [
                    {"clientId": int(client_id)},
                    {"propertyId": int(property_id)}
                ]
            },
            n_results=5,  # Fetch up to 5 relevant documents
            include=["documents", "metadatas"]
        )

        if not results.get("documents") or not results["documents"][0]:
            return jsonify(error="No embeddings found for given Client ID and Property ID"), 404

        # Format response with file names and content snippets
        docs = [
            {
                "fileName": meta["fileName"],
                "contentSnippet": doc[:500]  # Show first 500 characters
            }
            for doc, meta in zip(results["documents"][0], results["metadatas"][0])
        ]

        return jsonify(docs)

    except Exception as e:
        return jsonify(error=str(e)), 500


@app.route('/update-file-indexing', methods=['POST'])
def update_index():
    try:
        data = request.json
        required = ['clientId', 'indexedId', 'fileName', 'fileUrl', 'propertyId']
        if not all(k in data for k in required):
            return jsonify(error="Missing required fields"), 400
            
        content = get_file_content(data['fileUrl'])
        embedding = get_embedding(content)
        
        collection.update(
            ids=[data['indexedId']],
            embeddings=[embedding],
            metadatas=[{
                "clientId": data['clientId'],
                "propertyId": data['propertyId'],
                "fileName": data['fileName']
            }]
        )
        
        c.execute('''UPDATE file_indexes SET fileName=?, fileUrl=?, propertyId=?
                  WHERE clientId=? AND indexedId=?''',
                  (data['fileName'], data['fileUrl'], data['propertyId'],
                   data['clientId'], data['indexedId']))
        conn.commit()
        
        return jsonify({"indexed": True, "indexedid": data['indexedId']})
    
    except Exception as e:
        return jsonify(error=str(e)), 500

@app.route('/delete-indexing', methods=['POST'])
def delete_index():
    try:
        data = request.json
        required = ['clientId']
        if 'indexedid' not in data and 'indexedId' not in data:
            return jsonify(error="Missing index ID field"), 400
        if 'clientId' not in data:
            return jsonify(error="Missing client ID field"), 400
            
        indexed_id = data.get('indexedid') or data.get('indexedId')
        
        collection.delete(ids=[indexed_id])
        c.execute("DELETE FROM file_indexes WHERE clientId=? AND indexedId=?",
                  (data['clientId'], indexed_id))
        conn.commit()
        return jsonify({
            "deleted": True,
            "indexedid": indexed_id
        })
    
    except Exception as e:
        return jsonify(error=str(e)), 500
                
@app.route('/ask', methods=['POST'])
def ask():
    try:
        data = request.json
        required = ['clientId', 'query']  # Remove propertyId from required fields
        if not all(k in data for k in required):
            return jsonify(error="Missing required fields"), 400

        property_id = data.get('propertyId', None)  # Get propertyId, allow None

        # If propertyId is empty or not provided, use sample-global-faq.txt
        if not property_id:
            faq_file = "sample-global-faq.txt"
            if not os.path.exists(faq_file):
                return jsonify(error="Global FAQ file not found"), 404

            with open(faq_file, "r", encoding="utf-8") as f:
                faq_content = f.read()

            prompt = f"Context:\n{faq_content}\n\nQuestion: {data['query']}\nAnswer:"
            response = generate_response(prompt)
            return jsonify({"airesponse": response, "aisources": ["sample-global-faq.txt"]})

        # Convert property_id to int (if provided)
        property_id = int(property_id)

        query_embedding = get_embedding(data['query'])

        results = collection.query(
            query_embeddings=[query_embedding],
            where={
                "$and": [
                    {"clientId": int(data['clientId'])},
                    {"propertyId": property_id}
                ]
            },
            n_results=3,
            include=["documents", "metadatas"]
        )

        if not results["documents"]:
            return jsonify(error="No matching documents found"), 404

        # Build context from actual stored documents
        context_parts = []
        for doc, meta in zip(results['documents'][0], results['metadatas'][0]):
            context_parts.append(f"From {meta['fileName']}:\n{doc}")

        context = "\n\n".join(context_parts)
        prompt = f"Context:\n{context}\n\nQuestion: {data['query']}\nAnswer:"

        response = generate_response(prompt)
        return jsonify({
            "airesponse": response,
            "aisources": [m['fileName'] for m in results['metadatas'][0]]
        })

    except Exception as e:
        return jsonify(error=str(e)), 500

@app.route('/ask-generic', methods=['POST'])
def ask_generic():
    try:
        data = request.json
        if 'query' not in data:
            return jsonify(error="Missing query field"), 400
            
        response = generate_response(data['query'])
        return jsonify({"airesponse": response})
    
    except Exception as e:
        return jsonify(error=str(e)), 500

@app.route('/')
def index():
    return render_template('index.html')

if __name__ == '__main__':
    app.run(debug=True, use_reloader=False)  # Add use_reloader=False